from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── カラーパレット ──────────────────────────────
NAVY   = RGBColor(0x1A, 0x2F, 0x5E)
BLUE   = RGBColor(0x2E, 0x75, 0xB6)
LBLUE  = RGBColor(0xD6, 0xE4, 0xF7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
GRAY   = RGBColor(0x44, 0x44, 0x44)
MGRAY  = RGBColor(0x88, 0x88, 0x88)
GREEN  = RGBColor(0x1E, 0x7A, 0x45)
LGREEN = RGBColor(0xE2, 0xF4, 0xE8)
RED    = RGBColor(0xC0, 0x00, 0x00)
LRED   = RGBColor(0xFC, 0xE8, 0xE8)
AMBER  = RGBColor(0xBF, 0x6C, 0x00)
LAMBER = RGBColor(0xFD, 0xF2, 0xDC)

FONT = "游ゴシック"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── ユーティリティ ────────────────────────────────
def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=BLUE, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw
    else: s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h, size=14, bold=False,
       color=GRAY, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = FONT
    return box

def tb_multi(slide, lines, l, t, w, h, size=13, color=GRAY, spacing=1.15):
    """lines: list of (indent, text, bold, color_override)"""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    from pptx.util import Pt as PT
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    for i, item in enumerate(lines):
        indent, text, bold, col = (item + (None,))[:4] if len(item) < 4 else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        if indent:
            p.level = indent
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = col if col else color
    return box

def pic(slide, path, l, t, w, h):
    slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

def header(slide, title, sub=""):
    rect(slide, 0, 0, 13.33, 1.05, fill=NAVY)
    rect(slide, 0, 1.05, 13.33, 0.04, fill=BLUE)
    tb(slide, title, 0.35, 0.1, 9.5, 0.65, size=26, bold=True, color=WHITE)
    if sub:
        tb(slide, sub, 0.35, 0.67, 9.5, 0.38, size=12, color=LBLUE)
    rect(slide, 0, 7.18, 13.33, 0.32, fill=NAVY)

def insight_panel(slide, title, items, l, t, w, h, accent=BLUE):
    """右サイドの解説パネル"""
    rect(slide, l, t, w, h, fill=LGRAY)
    rect(slide, l, t, 0.07, h, fill=accent)
    tb(slide, title, l+0.18, t+0.1, w-0.25, 0.4,
       size=13, bold=True, color=accent)
    lines = [(0, item, False, GRAY) for item in items]
    tb_multi(slide, lines, l+0.18, t+0.55, w-0.3,
             h-0.65, size=12, color=GRAY)

def kpi_box(slide, label, value, l, t, w=2.5, h=1.1,
            bg_col=LBLUE, val_col=NAVY):
    rect(slide, l, t, w, h, fill=bg_col)
    tb(slide, label, l+0.1, t+0.08, w-0.15, 0.4,
       size=11, color=MGRAY, bold=False)
    tb(slide, value, l+0.1, t+0.45, w-0.15, 0.6,
       size=20, bold=True, color=val_col)

def tbl(slide, headers, rows, l, t, w, h,
        hdr_fill=NAVY, even_fill=LGRAY, odd_fill=WHITE):
    cols = len(headers)
    shape = slide.shapes.add_table(
        len(rows)+1, cols,
        Inches(l), Inches(t), Inches(w), Inches(h))
    tbl_obj = shape.table
    cw = Inches(w / cols)
    for c in range(cols): tbl_obj.columns[c].width = cw

    def cell(r, c, text, bold=False, bg=None, fg=WHITE, sz=11):
        cl = tbl_obj.cell(r, c); cl.text = text
        p = cl.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(sz); run.font.bold = bold
        run.font.name = FONT; run.font.color.rgb = fg
        if bg: cl.fill.solid(); cl.fill.fore_color.rgb = bg

    for c, h_ in enumerate(headers):
        cell(0, c, h_, bold=True, bg=hdr_fill, fg=WHITE, sz=11)
    for r, row in enumerate(rows):
        bg_ = even_fill if r % 2 == 0 else odd_fill
        for c, v in enumerate(row):
            cell(r+1, c, str(v), bg=bg_, fg=GRAY, sz=11)


# =====================================================================
# Slide 1: タイトル
# =====================================================================
s = add_slide()
bg(s, NAVY)
rect(s, 0, 2.2, 13.33, 3.1, fill=BLUE)
rect(s, 0, 2.2, 0.25, 3.1, fill=RGBColor(0xFF, 0xC0, 0x00))
tb(s, "課題データ集計・分析レポート",
   0.5, 2.35, 12.5, 1.1, size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
tb(s, "課題2（個人別スコア推移）／ 課題3（所属別パフォーマンス）",
   0.5, 3.45, 12, 0.55, size=18, color=LBLUE)
tb(s, "2024年9月　マーケティング観点分析",
   0.5, 4.05, 10, 0.45, size=14, color=LBLUE)
rect(s, 0, 6.8, 13.33, 0.7, fill=RGBColor(0x0F, 0x1E, 0x42))
tb(s, "CONFIDENTIAL　／　社内資料",
   0.5, 6.85, 12, 0.45, size=12, color=MGRAY, align=PP_ALIGN.RIGHT)

# =====================================================================
# Slide 2: エグゼクティブサマリー
# =====================================================================
s = add_slide()
bg(s)
header(s, "エグゼクティブサマリー", "両課題の主要インサイト一覧")

kpi_box(s, "分析対象者数（課題3・会社）", "35名", 0.3, 1.2)
kpi_box(s, "全体平均スコア（課題3・会社）", "84.9点", 2.95, 1.2)
kpi_box(s, "分析対象者数（課題2・学校）", "5名", 5.6, 1.2)
kpi_box(s, "期間中の全体成長幅", "+9.4点", 8.25, 1.2)
kpi_box(s, "最大成長者", "田中一郎 +9.3点", 10.9, 1.2, val_col=GREEN)

findings = [
    (0, "【課題3：会社データ｜所属別パフォーマンス】", True, NAVY),
    (0, "▶  開発部門が平均88.9点・ハイパフォーマー率40%でトップ。管理部門も87.1点と安定しており、",  False, GRAY),
    (0, "    両部門のベストプラクティスを組織全体に横展開することで底上げが期待できる。", False, GRAY),
    (0, "▶  人事部門は平均79.3点（全体比-5.6点）・ハイパフォーマー0名。標準偏差わずか2.29と", False, GRAY),
    (0, "    全員が低水準に均質化しており、採用・組織力全体のボトルネックになっている。", False, GRAY),
    (0, "▶  営業部門は平均83点だが標準偏差7.48と最大のばらつき。個人依存の構造が顕著。", False, GRAY),
    (0, "", False, GRAY),
    (0, "【課題2：学校データ｜科目別スコア推移】", True, NAVY),
    (0, "▶  全体平均は5日間で80.3点→89.8点（+9.5点）と右肩上がり。研修効果が明確に確認できる。", False, GRAY),
    (0, "▶  田中一郎（+9.3点）が最大の伸びを示し、投資対効果が最も高い人材と判断できる。", False, GRAY),
    (0, "▶  高橋花子は初日から92点超の高水準だが後半に-2点の微減。天井効果・学習疲れの兆候あり。", False, GRAY),
    (0, "▶  佐藤次郎は国語65点・社会68点と特定科目に集中した弱点があり、個別フォローで底上げ可能。", False, GRAY),
]
tb_multi(s, findings, 0.3, 2.42, 12.7, 4.6, size=12)

# =====================================================================
# Slide 3: セクション区切り（課題3）
# =====================================================================
s = add_slide()
bg(s, NAVY)
rect(s, 0, 2.8, 13.33, 2.2, fill=BLUE)
rect(s, 0, 2.8, 0.3, 2.2, fill=RGBColor(0xFF, 0xC0, 0x00))
tb(s, "課題３", 0.6, 2.85, 12, 0.55, size=18, bold=False, color=LBLUE)
tb(s, "会社データ｜所属別パフォーマンス分析",
   0.6, 3.3, 12, 0.9, size=34, bold=True, color=WHITE)
tb(s, "部門ごとのスコア分布・強み・課題・リスクを可視化する",
   0.6, 4.25, 12, 0.55, size=15, color=LBLUE)
rect(s, 0, 6.8, 13.33, 0.7, fill=RGBColor(0x0F, 0x1E, 0x42))

# =====================================================================
# Slide 4: 課題3 ドーナツ
# =====================================================================
s = add_slide()
bg(s)
header(s, "課題3｜所属別参加者数（会社データ）", "ドーナツ型円グラフ")
pic(s, "graph1_pie.png", 0.3, 1.15, 8.0, 5.85)
insight_panel(s,
    "グラフの読み方・ポイント",
    [
        "■ 参加者構成",
        "  営業・開発がそれぞれ10名（各29%）",
        "  管理8名（23%）、人事7名（20%）",
        "",
        "■ 注目点",
        "  営業・開発で全体の約6割を占める。",
        "  この2部門のパフォーマンスが",
        "  組織全体の数値を大きく左右する。",
        "",
        "■ マーケティング示唆",
        "  人事は参加者数も最少（7名）。",
        "  少人数で均質な低スコアという",
        "  構造は、集団研修より個別介入が",
        "  効果的であることを示唆する。",
    ],
    l=8.5, t=1.15, w=4.55, h=5.85, accent=BLUE)

# =====================================================================
# Slide 4: 課題3 グループ棒グラフ
# =====================================================================
s = add_slide()
bg(s)
header(s, "課題3｜所属別スコア比較（会社データ）", "平均・最高・最低 グループ棒グラフ")
pic(s, "graph2_bar.png", 0.3, 1.15, 8.0, 5.85)
insight_panel(s,
    "スコア比較の着眼点",
    [
        "■ 部門別ランキング（平均）",
        "  1位 開発   88.9点",
        "  2位 管理   87.1点",
        "  3位 営業   83.0点",
        "  4位 人事   79.3点",
        "",
        "■ 安定性（最高〜最低の差）",
        "  開発：94-82 ＝ 12点差（安定）",
        "  営業：93-70 ＝ 23点差（不安定）",
        "  人事：82-75 ＝  7点差（均質低）",
        "",
        "■ 読み取れること",
        "  開発は高水準かつ安定。一方、営業は",
        "  最高93点のトップがいても最低70点が",
        "  引き下げており「2極化」が進んでいる。",
        "  人事は差が小さいが水準そのものが低い。",
    ],
    l=8.5, t=1.15, w=4.55, h=5.85, accent=BLUE)

# =====================================================================
# Slide 5: 課題3 ヒストグラム
# =====================================================================
s = add_slide()
bg(s)
header(s, "課題3｜スコア分布ヒストグラム（会社データ）", "所属別色分け＋全体平均線（84.9点）")
pic(s, "graph3_histogram.png", 0.3, 1.15, 8.0, 5.85)
insight_panel(s,
    "分布形状から読む組織状態",
    [
        "■ 全体平均：84.9点",
        "  中央値86点 ＞ 平均84.9点",
        "  → 低スコア側に引っ張られる分布",
        "",
        "■ 右寄り（高得点）に集中する部門",
        "  開発・管理は90点台に多くの人材",
        "  → 高水準で安定したパフォーマンス",
        "",
        "■ 左寄り（低得点）が目立つ部門",
        "  人事は75〜82点帯に全員が集中",
        "  営業は70〜80点帯にも複数名",
        "",
        "■ マーケティング的解釈",
        "  人事の分布は「正規分布」ではなく",
        "  「圧縮された低水準の山」。",
        "  外部刺激（研修・人材交流）なしに",
        "  自然改善は期待しにくい状態。",
    ],
    l=8.5, t=1.15, w=4.55, h=5.85, accent=BLUE)

# =====================================================================
# Slide 6: 課題3 分析サマリー
# =====================================================================
s = add_slide()
bg(s)
header(s, "課題3｜マーケティング分析サマリー（会社データ）", "強み・課題・リスクの三層構造")

for i, (bx, title, accent, bg_c, items) in enumerate([
    (0.3, "強み（開発・管理）", GREEN, LGREEN, [
        "開発：平均88.9点・標準偏差3.57",
        "　→ 高品質かつ安定したアウトプット",
        "管理：平均87.1点・標準偏差5.62",
        "　→ 均整のとれた高水準チーム",
        "ハイパフォーマー率（90点以上）",
        "　開発40% / 管理37.5%",
        "",
        "【施策】",
        "成功パターンをマニュアル化し",
        "営業・人事へ横展開する。",
        "高スコア者をメンターに任命。",
    ]),
    (4.7, "課題（人事）", AMBER, LAMBER, [
        "平均79.3点（全体比 −5.6点）",
        "　→ 全部門で唯一80点未満",
        "ハイパフォーマー：0名",
        "標準偏差2.29（最小）",
        "　→ 全員が同じ低水準に均質化",
        "",
        "【なぜ問題か】",
        "人事が弱いと採用品質・",
        "組織施策・人材育成に連鎖影響。",
        "組織全体の底を規定する部門。",
        "",
        "【施策】外部研修・個別コーチング優先",
    ]),
    (9.1, "リスク（営業）", RED, LRED, [
        "平均83.0点・標準偏差7.48（最大）",
        "最低70点 ↔ 最高93点（差23点）",
        "ローパフォーマー（80点未満）4名",
        "　→ 全部門で最多",
        "",
        "【なぜ問題か】",
        "営業力の個人依存＝",
        "売上の予測可能性が低下。",
        "トップが抜けた瞬間に",
        "収益が急落するリスクがある。",
        "",
        "【施策】トップ営業のナレッジ移転・OJT",
    ]),
]):
    rect(s, bx, 1.15, 3.85, 6.0, fill=bg_c)
    rect(s, bx, 1.15, 3.85, 0.5, fill=accent)
    tb(s, title, bx+0.12, 1.2, 3.6, 0.42,
       size=14, bold=True, color=WHITE)
    lines = [(0, item, item.startswith("【"), accent if item.startswith("【") else GRAY)
             for item in items]
    tb_multi(s, lines, bx+0.15, 1.72, 3.6, 5.35, size=12)

# =====================================================================
# Slide 8: セクション区切り（課題2）
# =====================================================================
s = add_slide()
bg(s, RGBColor(0x1A, 0x3A, 0x2A))
rect(s, 0, 2.8, 13.33, 2.2, fill=RGBColor(0x1E, 0x7A, 0x45))
rect(s, 0, 2.8, 0.3, 2.2, fill=RGBColor(0xFF, 0xC0, 0x00))
tb(s, "課題２", 0.6, 2.85, 12, 0.55, size=18, bold=False, color=LGREEN)
tb(s, "学校データ｜科目別スコア分析",
   0.6, 3.3, 12, 0.9, size=34, bold=True, color=WHITE)
tb(s, "受講者ごとの科目別スコア推移・成長幅・弱点科目を可視化する",
   0.6, 4.25, 12, 0.55, size=15, color=LGREEN)
rect(s, 0, 6.8, 13.33, 0.7, fill=RGBColor(0x0D, 0x26, 0x18))

# =====================================================================
# Slide 9: 課題2 トレンド
# =====================================================================
s = add_slide()
bg(s)
header(s, "課題2｜日別スコア成長トレンド（学校データ）", "受講者別折れ線グラフ　／　全体平均（破線）")
pic(s, "graph4_trend.png", 0.3, 1.15, 8.0, 5.85)
insight_panel(s,
    "成長トレンドの読み方",
    [
        "■ 全体（黒破線）の動き",
        "  9/1：80.3点 → 9/5：89.8点",
        "  5日間で +9.5点の明確な成長",
        "  → 研修・学習効果が数値で確認できる",
        "",
        "■ 注目すべき個人",
        "  田中一郎（緑）",
        "  　初日81.7 → 最終日91.0（+9.3点）",
        "  　中盤以降に急加速するタイプ",
        "",
        "  高橋花子（赤）",
        "  　初日92.5 → 最終日90.5（−2.0点）",
        "  　序盤から高水準だが後半に失速",
        "",
        "■ 示唆",
        "  高初期スコア者は「維持」より",
        "  「次のレベル」への刺激が必要。",
        "  成長余地があるのは田中・佐藤。",
    ],
    l=8.5, t=1.15, w=4.55, h=5.85, accent=BLUE)

# =====================================================================
# Slide 8: 課題2 ヒートマップ
# =====================================================================
s = add_slide()
bg(s)
header(s, "課題2｜受講者×科目 スコアヒートマップ（学校データ）", "赤（低）→黄（中）→緑（高）のグラデーション")
pic(s, "graph5_heatmap.png", 0.3, 1.15, 8.0, 5.85)
insight_panel(s,
    "科目別・個人別の弱点マップ",
    [
        "■ 科目別平均（高→低）",
        "  英語   86.7点  ← 最も得意",
        "  社会   85.0点",
        "  数学   84.8点",
        "  理科   84.4点",
        "  国語   83.7点  ← 最も苦手",
        "",
        "■ 個人の弱点科目",
        "  佐藤次郎：社会68.0点・国語71.0点",
        "  　→ 文系科目に集中した弱点",
        "  山田太郎：国語81.0点が他より低い",
        "  鈴木花子：数学83.0点が他より低い",
        "",
        "■ マーケティング示唆",
        "  弱点は「人×科目」の組み合わせで",
        "  特定できる。全員一律の研修より",
        "  弱点科目への個別フォローが",
        "  最短で全体平均を引き上げる。",
    ],
    l=8.5, t=1.15, w=4.55, h=5.85, accent=BLUE)

# =====================================================================
# Slide 9: 課題2 箱ひげ＋成長幅
# =====================================================================
s = add_slide()
bg(s)
header(s, "課題2｜スコア分布・成長幅ランキング（学校データ）", "箱ひげ図（分布）＋横棒グラフ（初日→最終日の変化）")
pic(s, "graph6_boxplot_growth.png", 0.3, 1.15, 12.7, 5.85)

# =====================================================================
# Slide 10: 課題2 分析サマリー（テーブル）
# =====================================================================
s = add_slide()
bg(s)
header(s, "課題2｜受講者別パフォーマンス詳細（学校データ）", "平均・成長幅・弱点・推奨施策")

tbl(s,
    ["名前", "平均", "成長幅", "強み科目", "弱点科目", "総合評価・推奨施策"],
    [
        ["高橋花子", "92.0点", "▼ −2.0点", "理科96点", "後半失速",
         "トップ層。発展コンテンツ・上位資格挑戦で天井突破を狙う"],
        ["鈴木花子",  "88.6点", "▲ +4.7点", "英語92点", "理科80点",
         "安定した高水準。理科の強化で全科目90点台が射程圏内"],
        ["田中一郎",  "87.5点", "▲ +9.3点", "社会90点", "数学83点",
         "最大成長率。継続投資のROI最大。数学底上げで90点台狙い"],
        ["山田太郎",  "83.6点", "▲ +3.0点", "国語91点", "国語75点※",
         "ばらつきが課題。得意科目の安定化と苦手科目の反復練習を"],
        ["佐藤次郎",  "73.1点", "▲ +6.0点", "英語80点", "国語65点",
         "国語・社会が弱点の根本。文系科目への個別集中フォロー必須"],
    ],
    l=0.3, t=1.2, w=12.7, h=3.5)

tb(s, "※ 山田太郎の国語は初回75点→後半91点まで改善。反復学習の効果が最も顕著に現れた科目。",
   0.3, 4.75, 12.7, 0.4, size=11, color=MGRAY, italic=True)

tb_multi(s, [
    (0, "▶  田中一郎の「成長率+9.3点」は単なる数値ではなく、学習への適性が高い人材であることを示す。継続的なプログラムへの参加で組織の次世代リーダー候補として育成できる。", False, GRAY),
    (0, "▶  佐藤次郎の国語65点・社会68点は他受講者と比べて突出して低く、個別フォローなしの改善は見込みにくい。グループ学習に加え1on1セッションを並走させることを推奨する。", False, GRAY),
], 0.3, 5.2, 12.7, 1.85, size=12)

# =====================================================================
# Slide 11: 統合インサイト・推奨アクション
# =====================================================================
s = add_slide()
bg(s)
header(s, "統合インサイト｜優先推奨アクション", "課題2・課題3 マーケティング観点 総括")

actions = [
    (RED,   "★★★", "人事部門への集中支援（最優先）",
     "全体平均比 −5.6点。7名全員が低水準に均質化しており自然回復は困難。外部研修・個別コーチングを即時投入。"
     "人事の改善は採用品質→組織力→全部門のパフォーマンスに連鎖する。"),
    (RED,   "★★★", "営業ローパフォーマー4名の底上げ（最優先）",
     "ローパフォーマーが全部門最多で、最低70点と水準も低い。トップ営業（93点）のナレッジを"
     "OJT・ロープレで移転。営業力の組織化が直接的な収益安定化につながる。"),
    (BLUE,  "★★",  "田中一郎への継続投資",
     "成長幅+9.3点は全受講者最大。学習適性が高く、今後のプログラム継続でさらなる成長が期待できる。"
     "組織の次世代リーダー候補として育成ロードマップを設計する。"),
    (BLUE,  "★★",  "高橋花子に発展コンテンツを提供",
     "初日から92点超と突出した実力を持つが、後半に−2点の微減。天井効果・慣れによる失速を防ぐため、"
     "難易度の高い応用課題や外部認定試験への挑戦で学習意欲を維持する。"),
    (GRAY,  "★",   "開発・管理のナレッジを全社横展開",
     "高スコア・低ばらつきの両部門のメソッドをマニュアル・社内勉強会として可視化。"
     "特に営業・人事への水平展開が全体平均を84.9点→90点へ押し上げる近道となる。"),
]

for i, (accent, star, title, body) in enumerate(actions):
    by = 1.2 + i * 1.2
    rect(s, 0.3, by, 0.55, 1.05, fill=accent)
    tb(s, star, 0.3, by+0.28, 0.55, 0.5,
       size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    col = LGRAY if i % 2 == 0 else WHITE
    rect(s, 0.9, by, 12.1, 1.05, fill=col)
    tb(s, title, 1.05, by+0.06, 11.8, 0.4, size=13, bold=True, color=NAVY)
    tb(s, body,  1.05, by+0.48, 11.8, 0.55, size=11, color=GRAY)

# =====================================================================
# 保存
# =====================================================================
out = "課題データ分析レポート.pptx"
prs.save(out)
print(f"保存完了: {out}")
print(f"スライド数: {len(prs.slides)}枚")
